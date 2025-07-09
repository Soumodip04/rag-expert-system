"""Document processing pipeline for various file formats"""
import os
from typing import List, Dict, Any, Optional
from pathlib import Path
from dataclasses import dataclass
from loguru import logger
import mimetypes

@dataclass
class ProcessedDocument:
    """Represents a processed document"""
    content: str
    metadata: Dict[str, Any] = None  # Use None as default, we'll initialize in __post_init__
    source_path: str = ""
    doc_type: str = "unknown"
    
    def __post_init__(self):
        """Ensure metadata is always a dictionary"""
        if self.metadata is None:
            self.metadata = {}
        elif not isinstance(self.metadata, dict):
            # Handle non-dict metadata by storing it in a special key
            original = self.metadata
            self.metadata = {"_original_metadata": original}
            # Log this conversion for debugging
            logger.debug(f"Converted non-dict metadata to dict: {type(original).__name__}")

class DocumentProcessor:
    """Base class for document processors"""
    
    def __init__(self):
        self.supported_formats = []
    
    def can_process(self, file_path: str) -> bool:
        """Check if this processor can handle the file"""
        return Path(file_path).suffix.lower() in self.supported_formats
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process a document and return structured content"""
        raise NotImplementedError

class PDFProcessor(DocumentProcessor):
    """PDF document processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.pdf']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process PDF file"""
        try:
            import PyPDF2
            
            content = ""
            doc_metadata = metadata or {}
            
            with open(file_path, 'rb') as file:
                try:
                    pdf_reader = PyPDF2.PdfReader(file)
                    
                    # Extract metadata from PDF
                    if pdf_reader.metadata:
                        doc_metadata.update({
                            'title': pdf_reader.metadata.get('/Title', ''),
                            'author': pdf_reader.metadata.get('/Author', ''),
                            'subject': pdf_reader.metadata.get('/Subject', ''),
                            'creator': pdf_reader.metadata.get('/Creator', ''),
                            'pages': len(pdf_reader.pages)
                        })
                    
                    # Extract text from all pages with better error handling
                    for page_num, page in enumerate(pdf_reader.pages):
                        try:
                            # Handle encoding issues
                            page_text = None
                            # Try with different extraction methods
                            try:
                                page_text = page.extract_text()
                            except Exception as e:
                                logger.warning(f"Standard extraction failed on page {page_num + 1}: {e}")
                                try:
                                    # Try layout extraction
                                    page_text = page.extract_text(extraction_mode="layout")
                                except Exception as e2:
                                    logger.warning(f"Layout extraction failed on page {page_num + 1}: {e2}")
                                    try:
                                        # Try raw extraction
                                        page_text = page.extract_text(extraction_mode="raw")
                                    except Exception as e3:
                                        logger.warning(f"Raw extraction failed on page {page_num + 1}: {e3}")
                            
                            # If we got text, sanitize and add it
                            if page_text and page_text.strip():
                                # Replace problematic characters
                                page_text = ''.join(char if ord(char) < 65536 else ' ' for char in page_text)
                                content += f"\n\n--- Page {page_num + 1} ---\n{page_text}"
                            else:
                                # Try direct extraction by character ranges
                                content += f"\n\n--- Page {page_num + 1} (Limited Text) ---\n"
                                try:
                                    # Try character by character extraction
                                    text_pieces = []
                                    for i in range(10):  # Try 10 sections
                                        try:
                                            start = i * 500
                                            chunk = page.extract_text(start, start + 500)
                                            if chunk:
                                                # Remove problematic characters
                                                chunk = ''.join(char if ord(char) < 65536 else ' ' for char in chunk)
                                                text_pieces.append(chunk)
                                        except:
                                            pass
                                    if text_pieces:
                                        content += " ".join(text_pieces)
                                except Exception as chunk_err:
                                    logger.warning(f"Chunk extraction failed on page {page_num + 1}: {chunk_err}")
                        except Exception as e:
                            logger.warning(f"Error extracting text from page {page_num + 1}: {e}")
                            # Add page marker for tracking
                            content += f"\n\n--- Page {page_num + 1} (Extraction Failed) ---\n"
                except Exception as pdf_err:
                    logger.error(f"Error reading PDF with PyPDF2: {pdf_err}")
                    # Try alternate PDF reading method
                    content = self._fallback_pdf_extract(file_path)
            
            # If we have no content, try the fallback method
            if not content.strip():
                content = self._fallback_pdf_extract(file_path)
            
            doc_metadata['file_size'] = os.path.getsize(file_path)
            doc_metadata['file_name'] = Path(file_path).name
            
            return ProcessedDocument(
                content=content.strip() if content.strip() else "[PDF content extraction failed - no text content found]",
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="pdf"
            )
        
        except ImportError:
            logger.error("PyPDF2 not installed. Please install with: pip install PyPDF2")
            raise

    def _fallback_pdf_extract(self, file_path: str) -> str:
        """Fallback method for PDF text extraction when PyPDF2 fails"""
        # First attempt: Try with pdfplumber if available
        try:
            import pdfplumber
            content = ""
            
            with pdfplumber.open(file_path) as pdf:
                for i, page in enumerate(pdf.pages):
                    try:
                        text = page.extract_text()
                        if text:
                            # Clean problematic Unicode characters
                            text = ''.join(char if ord(char) < 65536 else ' ' for char in text)
                            content += f"\n\n--- Page {i+1} ---\n{text}"
                    except Exception as e:
                        logger.warning(f"pdfplumber extraction error on page {i+1}: {e}")
            
            if content.strip():
                logger.info("Used pdfplumber fallback for PDF extraction")
                return content
        except ImportError:
            logger.warning("pdfplumber not available for fallback PDF extraction")
        except Exception as e:
            logger.warning(f"pdfplumber failed: {e}")
        
        # Second attempt: Try PyPDF2 again with different approach
        try:
            import PyPDF2
            content = ""
            
            with open(file_path, 'rb') as file:
                try:
                    pdf_reader = PyPDF2.PdfReader(file)
                    
                    # Try a different approach - extract raw data and sanitize
                    for page_num in range(len(pdf_reader.pages)):
                        try:
                            page = pdf_reader.pages[page_num]
                            # Try extraction modes in sequence
                            extraction_modes = ["layout", "plain", "raw", "structured"]
                            for mode in extraction_modes:
                                try:
                                    text = page.extract_text(extraction_mode=mode) if hasattr(page, 'extract_text') else page.extract_text()
                                    if text and text.strip():
                                        # Clean problematic Unicode characters
                                        text = ''.join(char if ord(char) < 65536 else ' ' for char in text)
                                        content += f"\n\n--- Page {page_num + 1} ---\n{text}"
                                        break
                                except Exception:
                                    continue
                                    
                            if not text or not text.strip():
                                logger.warning(f"All PyPDF2 extraction modes failed for page {page_num + 1}")
                        except Exception as e:
                            logger.warning(f"PyPDF2 second attempt extraction error on page {page_num + 1}: {e}")
                except Exception as e:
                    logger.warning(f"PyPDF2 second attempt failed: {e}")
            
            if content.strip():
                logger.info("Used PyPDF2 alternative extraction method successfully")
                return content
        except Exception as e:
            logger.warning(f"PyPDF2 alternative extraction failed: {e}")
            
        # Third attempt: Raw binary extraction as a fallback
        try:
            # Read file as binary and look for text streams
            with open(file_path, 'rb') as f:
                pdf_data = f.read()
                
            # Basic text extraction using string operations
            content = ""
            
            # Extract printable ASCII characters
            import re
            # More inclusive pattern to capture text in PDFs
            printable = re.compile(b'[A-Za-z0-9 .,;:!?\'"\\-+*/=<>()\\[\\]{}|#$%&@^_`~\n\r\t]{4,}')
            matches = printable.findall(pdf_data)
            
            text_parts = []
            for match in matches:
                if len(match) > 4:  # Ignore very short matches
                    # Try different encodings
                    encodings = ['utf-8', 'latin-1', 'ascii', 'windows-1252']
                    for encoding in encodings:
                        try:
                            text_part = match.decode(encoding, errors='replace')
                            # If it contains substantial text (not just symbols or garbage)
                            if re.search('[A-Za-z]{3,}', text_part):
                                text_parts.append(text_part)
                                break
                        except:
                            continue
            
            if text_parts:
                # Filter out duplicates and very short fragments
                unique_parts = []
                for part in text_parts:
                    if part not in unique_parts and len(part) > 10:
                        unique_parts.append(part)
                
                content = "\n".join(unique_parts)
                logger.info("Used raw binary extraction method for PDF")
                return content
                
        except Exception as e:
            logger.warning(f"Raw binary extraction failed: {e}")
        
        # OCR approach removed as it's dependent on external libraries that might not be installed
        
        # Last resort: Extract metadata only
        try:
            import PyPDF2
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                text = "[Limited PDF content extracted]\n\n"
                text += f"Total pages: {len(pdf_reader.pages)}\n"
                
                # Try to get at least some metadata
                if pdf_reader.metadata:
                    text += "Document metadata:\n"
                    for key, value in pdf_reader.metadata.items():
                        if value:
                            try:
                                safe_value = str(value).encode('ascii', 'replace').decode()
                                text += f"{key}: {safe_value}\n"
                            except:
                                pass
                            
                return text
        except Exception:
            pass
            
        # If absolutely everything fails
        logger.error("All PDF extraction methods failed")
        return "[PDF content could not be extracted - please check file format]"

class DocxProcessor(DocumentProcessor):
    """DOCX document processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.docx', '.doc']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process DOCX file"""
        try:
            import docx2txt
            
            content = docx2txt.process(file_path)
            doc_metadata = metadata or {}
            
            # Add basic file metadata
            doc_metadata.update({
                'file_size': os.path.getsize(file_path),
                'file_name': Path(file_path).name,
                'word_count': len(content.split()) if content else 0
            })
            
            return ProcessedDocument(
                content=content or "",
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="docx"
            )
        
        except ImportError:
            logger.error("docx2txt not installed. Please install with: pip install docx2txt")
            raise
        except Exception as e:
            logger.error(f"Error processing DOCX {file_path}: {e}")
            return ProcessedDocument(
                content=f"Error processing DOCX: {str(e)}",
                metadata=metadata or {},
                source_path=file_path,
                doc_type="docx"
            )


class TextProcessor(DocumentProcessor):
    """Plain text document processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.txt', '.md', '.rst', '.log']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process text file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            content = ""
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    break
                except UnicodeDecodeError:
                    continue
            
            if not content:
                logger.warning(f"Could not decode file {file_path} with any supported encoding")
            
            doc_metadata = metadata or {}
            doc_metadata.update({
                'file_size': os.path.getsize(file_path),
                'file_name': Path(file_path).name,
                'line_count': len(content.splitlines()) if content else 0,
                'word_count': len(content.split()) if content else 0
            })
            
            return ProcessedDocument(
                content=content,
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="text"
            )
        
        except Exception as e:
            logger.error(f"Error processing text file {file_path}: {e}")
            return ProcessedDocument(
                content=f"Error processing text file: {str(e)}",
                metadata=metadata or {},
                source_path=file_path,
                doc_type="text"
            )


class HTMLProcessor(DocumentProcessor):
    """HTML document processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.html', '.htm']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process HTML file"""
        try:
            from bs4 import BeautifulSoup
            
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Extract metadata from HTML
            doc_metadata = metadata or {}
            
            title_tag = soup.find('title')
            if title_tag:
                doc_metadata['title'] = title_tag.get_text().strip()
            
            # Extract meta tags
            meta_tags = soup.find_all('meta')
            for tag in meta_tags:
                if tag.get('name') and tag.get('content'):
                    doc_metadata[f"meta_{tag.get('name')}"] = tag.get('content')
            
            # Extract text content
            # Remove script and style elements
            for script in soup(["script", "style"]):
                script.decompose()
            
            content = soup.get_text()
            
            # Clean up text
            lines = (line.strip() for line in content.splitlines())
            chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
            content = ' '.join(chunk for chunk in chunks if chunk)
            
            doc_metadata.update({
                'file_size': os.path.getsize(file_path),
                'file_name': Path(file_path).name,
                'word_count': len(content.split()) if content else 0
            })
            
            return ProcessedDocument(
                content=content,
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="html"
            )
        
        except ImportError:
            logger.error("BeautifulSoup4 not installed. Please install with: pip install beautifulsoup4")
            raise
        except Exception as e:
            logger.error(f"Error processing HTML {file_path}: {e}")
            return ProcessedDocument(
                content=f"Error processing HTML: {str(e)}",
                metadata=metadata or {},
                source_path=file_path,
                doc_type="html"
            )


class CSVProcessor(DocumentProcessor):
    """CSV document processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.csv']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process CSV file"""
        try:
            import pandas as pd
            
            # Read CSV
            df = pd.read_csv(file_path)
            
            # Convert to text representation
            content = f"CSV Data Summary:\n"
            content += f"Rows: {len(df)}, Columns: {len(df.columns)}\n\n"
            content += f"Column Names: {', '.join(df.columns.tolist())}\n\n"
            
            # Add first few rows as sample
            content += "Sample Data:\n"
            content += df.head().to_string(index=False)
            
            # Add column statistics for numeric columns
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                content += "\n\nNumeric Column Statistics:\n"
                content += df[numeric_cols].describe().to_string()
            
            doc_metadata = metadata or {}
            doc_metadata.update({
                'file_size': os.path.getsize(file_path),
                'file_name': Path(file_path).name,
                'rows': len(df),
                'columns': len(df.columns),
                'column_names': df.columns.tolist()
            })
            
            return ProcessedDocument(
                content=content,
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="csv"
            )
        
        except ImportError:
            logger.error("pandas not installed. Please install with: pip install pandas")
            raise
        except Exception as e:
            logger.error(f"Error processing CSV {file_path}: {e}")
            return ProcessedDocument(
                content=f"Error processing CSV: {str(e)}",
                metadata=metadata or {},
                source_path=file_path,
                doc_type="csv"
            )


class JSONProcessor(DocumentProcessor):
    """JSON document processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.json']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process JSON file"""
        try:
            import json
            
            with open(file_path, 'r', encoding='utf-8') as file:
                data = json.load(file)
            
            # Convert JSON to readable text
            content = f"JSON Document Content:\n\n"
            content += json.dumps(data, indent=2, ensure_ascii=False)
            
            doc_metadata = metadata or {}
            doc_metadata.update({
                'file_size': os.path.getsize(file_path),
                'file_name': Path(file_path).name,
                'json_keys': list(data.keys()) if isinstance(data, dict) else [],
                'data_type': type(data).__name__
            })
            
            return ProcessedDocument(
                content=content,
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="json"
            )
        
        except Exception as e:
            logger.error(f"Error processing JSON {file_path}: {e}")
            return ProcessedDocument(
                content=f"Error processing JSON: {str(e)}",
                metadata=metadata or {},
                source_path=file_path,
                doc_type="json"
            )

class ExcelProcessor(DocumentProcessor):
    """Excel document processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.xlsx', '.xls']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process Excel file"""
        try:
            import pandas as pd
            
            # Read all sheets
            excel_file = pd.ExcelFile(file_path)
            content = f"Excel Document: {Path(file_path).name}\n\n"
            
            for sheet_name in excel_file.sheet_names:
                df = pd.read_excel(file_path, sheet_name=sheet_name)
                content += f"Sheet: {sheet_name}\n"
                content += f"Rows: {len(df)}, Columns: {len(df.columns)}\n"
                content += f"Columns: {', '.join(df.columns.tolist())}\n\n"
                
                # Add sample data
                if not df.empty:
                    content += "Sample Data:\n"
                    content += df.head(3).to_string(index=False)
                    content += "\n\n"
            
            doc_metadata = metadata or {}
            doc_metadata.update({
                'file_size': os.path.getsize(file_path),
                'file_name': Path(file_path).name,
                'sheet_names': excel_file.sheet_names,
                'sheet_count': len(excel_file.sheet_names)
            })
            
            return ProcessedDocument(
                content=content,
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="excel"
            )
        
        except Exception as e:
            logger.error(f"Error processing Excel {file_path}: {e}")
            return ProcessedDocument(
                content=f"Error processing Excel: {str(e)}",
                metadata=metadata or {},
                source_path=file_path,
                doc_type="excel"
            )

class PowerPointProcessor(DocumentProcessor):
    """PowerPoint document processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.pptx', '.ppt']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process PowerPoint file"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            content = f"PowerPoint Presentation: {Path(file_path).name}\n\n"
            
            for i, slide in enumerate(prs.slides, 1):
                content += f"Slide {i}:\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content += f"{shape.text.strip()}\n"
                
                content += "\n"
            
            doc_metadata = metadata or {}
            doc_metadata.update({
                'file_size': os.path.getsize(file_path),
                'file_name': Path(file_path).name,
                'slide_count': len(prs.slides)
            })
            
            return ProcessedDocument(
                content=content,
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="powerpoint"
            )
        
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {e}")
            return ProcessedDocument(
                content=f"Error processing PowerPoint: {str(e)}",
                metadata=metadata or {},
                source_path=file_path,
                doc_type="powerpoint"
            )

class YAMLProcessor(DocumentProcessor):
    """YAML document processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.yaml', '.yml']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process YAML file"""
        try:
            import yaml
            
            with open(file_path, 'r', encoding='utf-8') as file:
                data = yaml.safe_load(file)
            
            # Convert YAML to readable text
            content = f"YAML Document Content:\n\n"
            content += yaml.dump(data, default_flow_style=False, allow_unicode=True)
            
            doc_metadata = metadata or {}
            doc_metadata.update({
                'file_size': os.path.getsize(file_path),
                'file_name': Path(file_path).name,
                'yaml_keys': list(data.keys()) if isinstance(data, dict) else [],
                'data_type': type(data).__name__
            })
            
            return ProcessedDocument(
                content=content,
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="yaml"
            )
        
        except Exception as e:
            logger.error(f"Error processing YAML {file_path}: {e}")
            return ProcessedDocument(
                content=f"Error processing YAML: {str(e)}",
                metadata=metadata or {},
                source_path=file_path,
                doc_type="yaml"
            )

class CodeProcessor(DocumentProcessor):
    """Code file processor"""
    
    def __init__(self):
        super().__init__()
        self.supported_formats = ['.py', '.js', '.ts', '.java', '.cpp', '.c', '.h', '.css', '.sql']
    
    def process(self, file_path: str, metadata: Dict[str, Any] = None) -> ProcessedDocument:
        """Process code file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            # Add file type context
            file_ext = Path(file_path).suffix
            lang_map = {
                '.py': 'Python', '.js': 'JavaScript', '.ts': 'TypeScript',
                '.java': 'Java', '.cpp': 'C++', '.c': 'C', '.h': 'C/C++ Header',
                '.css': 'CSS', '.sql': 'SQL'
            }
            
            language = lang_map.get(file_ext, 'Code')
            formatted_content = f"{language} Source Code: {Path(file_path).name}\n\n"
            formatted_content += content
            
            doc_metadata = metadata or {}
            doc_metadata.update({
                'file_size': os.path.getsize(file_path),
                'file_name': Path(file_path).name,
                'language': language,
                'line_count': len(content.splitlines()),
                'file_extension': file_ext
            })
            
            return ProcessedDocument(
                content=formatted_content,
                metadata=doc_metadata,
                source_path=file_path,
                doc_type="code"
            )
        
        except Exception as e:
            logger.error(f"Error processing code file {file_path}: {e}")
            return ProcessedDocument(
                content=f"Error processing code file: {str(e)}",
                metadata=metadata or {},
                source_path=file_path,
                doc_type="code"
            )

class DocumentProcessingPipeline:
    """Document processing pipeline that handles multiple file formats"""
    
    def __init__(self):
        self.processors = [
            PDFProcessor(),
            DocxProcessor(),
            TextProcessor(),
            HTMLProcessor(),
            CSVProcessor(),
            JSONProcessor(),
            ExcelProcessor(),
            PowerPointProcessor(),
            YAMLProcessor(),
            CodeProcessor()
        ]
    
    def get_processor(self, file_path: str) -> Optional[DocumentProcessor]:
        """Get appropriate processor for file"""
        for processor in self.processors:
            if processor.can_process(file_path):
                return processor
        return None
    
    def classify_document_content(self, content: str) -> Dict[str, Any]:
        """Classify document content to add additional metadata
        
        This helps the model better understand the type and purpose of documents.
        """
        document_classification = {}
        
        # Check if content contains specific keywords for document type classification
        content_lower = content.lower()
        
        # Detect syllabus/curriculum documents
        if any(term in content_lower for term in ['syllabus', 'curriculum', 'course outline', 'learning objectives']):
            document_classification['doc_category'] = 'curriculum'
            document_classification['doc_purpose'] = 'educational'
            
        # Detect assignment/exam documents
        if any(term in content_lower for term in ['assignment', 'exam', 'test', 'quiz', 'question paper']):
            document_classification['doc_category'] = 'assessment'
            
            # Further classify assessment type
            if any(term in content_lower for term in ['exam', 'final exam', 'midterm']):
                document_classification['assessment_type'] = 'exam'
            elif any(term in content_lower for term in ['assignment', 'homework']):
                document_classification['assessment_type'] = 'assignment'
            elif any(term in content_lower for term in ['quiz', 'test']):
                document_classification['assessment_type'] = 'quiz'
                
        # Detect if document contains questions
        question_indicators = ['question', 'answer the following', 'solve the following', 'explain', 'describe']
        question_count = sum(content_lower.count(indicator) for indicator in question_indicators)
        if question_count > 3:  # Threshold to identify question documents
            document_classification['contains_questions'] = True
            document_classification['estimated_question_count'] = question_count
            
        # Detect technical content
        technical_terms = ['algorithm', 'function', 'programming', 'code', 'implementation', 'system', 'database', 'architecture']
        technical_score = sum(content_lower.count(term) for term in technical_terms)
        if technical_score > 5:
            document_classification['content_type'] = 'technical'
            
        # Detect legal content
        legal_terms = ['law', 'legal', 'rights', 'policy', 'agreement', 'terms', 'contract', 'regulation']
        legal_score = sum(content_lower.count(term) for term in legal_terms)
        if legal_score > 5:
            document_classification['content_type'] = 'legal'
            
        return document_classification
    
    def process_file(self, file_path: str, metadata: Dict[str, Any] = None) -> Optional[ProcessedDocument]:
        """Process a single file"""
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return None
        
        processor = self.get_processor(file_path)
        if not processor:
            logger.warning(f"No processor found for file: {file_path}")
            return None
        
        try:
            logger.info(f"Processing file: {file_path}")
            processed_doc = processor.process(file_path, metadata)
            
            # Add additional document classification
            if processed_doc:
                # Initialize metadata as dict if None or not a dict
                if processed_doc.metadata is None or not isinstance(processed_doc.metadata, dict):
                    processed_doc.metadata = {}
                    
                # Get classification results
                classification = self.classify_document_content(processed_doc.content)
                
                # Always make sure we have a dict to update
                if not processed_doc.metadata:
                    processed_doc.metadata = {}
                
                # Ensure classification is a dict before updating
                if isinstance(classification, dict):
                    processed_doc.metadata.update(classification)
                elif classification is not None:
                    # If classification is not a dict but has a value, store it in a special key
                    processed_doc.metadata["classification"] = classification
                    
                logger.info(f"Document classification for {file_path}: {classification}")
                
            return processed_doc
        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            return None
    
    def process_directory(self, directory_path: str, recursive: bool = True, 
                         metadata: Dict[str, Any] = None) -> List[ProcessedDocument]:
        """Process all supported files in a directory"""
        if not os.path.exists(directory_path):
            logger.error(f"Directory not found: {directory_path}")
            return []
        
        processed_docs = []
        
        if recursive:
            for root, dirs, files in os.walk(directory_path):
                for file in files:
                    file_path = os.path.join(root, file)
                    doc = self.process_file(file_path, metadata)
                    if doc:
                        processed_docs.append(doc)
        else:
            for file in os.listdir(directory_path):
                file_path = os.path.join(directory_path, file)
                if os.path.isfile(file_path):
                    doc = self.process_file(file_path, metadata)
                    if doc:
                        processed_docs.append(doc)
        
        logger.info(f"Processed {len(processed_docs)} documents from {directory_path}")
        return processed_docs
    
    def get_supported_formats(self) -> List[str]:
        """Get all supported file formats"""
        formats = []
        for processor in self.processors:
            formats.extend(processor.supported_formats)
        return sorted(list(set(formats)))
